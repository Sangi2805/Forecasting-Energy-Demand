"""Extract images from Draft.pptx for appendix slides."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

DRAFT = Path("/Users/msshakeel/Downloads/Draft.pptx")
OUT_DIR = Path(__file__).resolve().parent / "draft_images"

# Slide numbers (1-based) to keep in appendix — graph-heavy EDA detail
APPENDIX_SLIDES = [2, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15, 16, 17, 18]


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(DRAFT))
    for idx in APPENDIX_SLIDES:
        slide = prs.slides[idx - 1]
        img_i = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                ext = shape.image.ext
                blob = shape.image.blob
            except ValueError:
                continue
            path = OUT_DIR / f"slide{idx:02d}_{img_i}{ext}"
            path.write_bytes(blob)
            img_i += 1
        print(f"slide {idx}: {img_i} image(s)")


if __name__ == "__main__":
    main()

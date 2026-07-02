"""Build unified Capstone Iteration 1 deck — condensed main + EDA appendix."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ASSETS = Path(__file__).resolve().parent
DRAFT_IMAGES = ASSETS / "draft_images"
OUT = Path("/Users/msshakeel/Downloads/Capstone_Iteration1_Presentation.pptx")
OUT_PROJECT = ASSETS.parent / "Capstone_Iteration1_Presentation.pptx"

# ── Theme (navy + electric blue + warm accent) ───────────────────────────────
NAVY = RGBColor(0x0F, 0x29, 0x42)
BLUE = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_SOFT = RGBColor(0xE0, 0xE7, 0xFF)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF4, 0xF7, 0xFB)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Segoe UI"

APPENDIX_LABELS = {
    2: "Proposed timeline",
    4: "Data sources",
    5: "Preprocessing pipeline",
    6: "Demand over time",
    7: "Seasonal demand peaks",
    8: "Monthly demand peaks",
    9: "Long-term demand trend",
    10: "Demand patterns summary",
    11: "Hourly demand patterns",
    12: "Weekly demand patterns",
    13: "Seasonal & holiday demand",
    14: "Temperature vs demand",
    15: "Apparent temperature vs demand",
    16: "Snowfall vs demand",
    17: "Wind speed vs demand",
    18: "Correlation matrix",
}


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        self._blank = self.prs.slide_layouts[6] if len(self.prs.slide_layouts) > 6 else self.prs.slide_layouts[0]
        self._n = 0

    def _slide(self):
        self._n += 1
        slide = self.prs.slides.add_slide(self._blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
        return slide

    def _bar(self, slide, color: RGBColor = NAVY, height: float = 0.52) -> None:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(height))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, self.prs.slide_width - Inches(2.2), 0, Inches(2.2), Inches(height)
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = BLUE
        stripe.line.fill.background()

    def _footer(self, slide) -> None:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(5.38), Inches(8.9), Inches(0.015)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BLUE_SOFT
        line.line.fill.background()
        tag = slide.shapes.add_textbox(Inches(0.55), Inches(5.42), Inches(6.5), Inches(0.2))
        p = tag.text_frame.paragraphs[0]
        p.text = "Energy Demand Forecast · Capstone Iteration 1"
        p.font.name = FONT
        p.font.size = Pt(8)
        p.font.color.rgb = GRAY
        num = slide.shapes.add_textbox(Inches(9.0), Inches(5.42), Inches(0.45), Inches(0.2))
        np = num.text_frame.paragraphs[0]
        np.text = str(self._n)
        np.font.name = FONT
        np.font.size = Pt(8)
        np.font.color.rgb = GRAY
        np.alignment = PP_ALIGN.RIGHT

    def _panel(self, slide, left: float, top: float, width: float, height: float) -> None:
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = PANEL
        panel.line.color.rgb = BLUE_SOFT
        panel.line.width = Pt(1)

    def _accent(self, slide, top: float = 0.88) -> None:
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(top), Inches(0.85), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ORANGE
        line.line.fill.background()

    def _title(self, slide, text: str, top: float = 0.14, size: int = 26, on_bar: bool = True) -> None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.5), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.bold = True
        p.font.color.rgb = WHITE if on_bar else DARK

    def _subtitle(self, slide, text: str, top: float = 0.82) -> None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(8.9), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

    def _bullets(self, slide, items: list[str], top: float = 1.25, size: int = 16, width: float = 8.9) -> None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(top), Inches(width), Inches(4.2))
        tf = box.text_frame
        tf.word_wrap = True
        for i, text in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.name = FONT
            p.font.size = Pt(size)
            p.font.color.rgb = DARK
            p.space_after = Pt(8)
            p.level = 0

    def _caption(self, slide, text: str) -> None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(5.2), Inches(8.9), Inches(0.3))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    def title_slide(self, title: str, lines: list[str]) -> None:
        slide = self._slide()
        hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(3.2), self.prs.slide_height)
        hero.fill.solid()
        hero.fill.fore_color.rgb = NAVY
        hero.line.fill.background()
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.2), 0, Inches(0.12), self.prs.slide_height)
        band.fill.solid()
        band.fill.fore_color.rgb = ORANGE
        band.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(3.55), Inches(1.55), Inches(6.1), Inches(1.5))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = DARK

        meta = slide.shapes.add_textbox(Inches(3.55), Inches(3.2), Inches(6.1), Inches(1.8))
        tf = meta.text_frame
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = FONT
            p.font.size = Pt(12 if i > 0 else 14)
            p.font.color.rgb = GRAY if i > 0 else BLUE
            p.font.bold = i == 0
        self._footer(slide)

    def section(self, title: str, subtitle: str = "") -> None:
        slide = self._slide()
        block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        block.fill.solid()
        block.fill.fore_color.rgb = NAVY
        block.line.fill.background()
        orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(-0.8), Inches(3.5), Inches(3.5))
        orb.fill.solid()
        orb.fill.fore_color.rgb = BLUE
        orb.fill.transparency = 0.25
        orb.line.fill.background()

        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.85), Inches(1.45), Inches(0.34)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = ORANGE
        pill.line.fill.background()
        pp = pill.text_frame.paragraphs[0]
        pp.text = "SECTION"
        pp.font.name = FONT
        pp.font.size = Pt(9)
        pp.font.bold = True
        pp.font.color.rgb = WHITE
        pp.alignment = PP_ALIGN.CENTER
        pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        tb = slide.shapes.add_textbox(Inches(0.65), Inches(2.35), Inches(8.5), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = FONT
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(0.65), Inches(3.25), Inches(8.5), Inches(0.5))
            sp = sub.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.name = FONT
            sp.font.size = Pt(15)
            sp.font.color.rgb = BLUE_LIGHT
        self._footer(slide)

    def content(self, title: str, bullets: list[str], subtitle: str = "") -> None:
        slide = self._slide()
        self._bar(slide)
        self._title(slide, title)
        self._panel(slide, 0.45, 0.72, 9.1, 4.55)
        top = 1.35 if subtitle else 1.15
        if subtitle:
            self._subtitle(slide, subtitle, top=0.72)
        self._bullets(slide, bullets, top=top, size=15)
        self._footer(slide)

    def two_col(self, title: str, left: list[str], right: list[str]) -> None:
        slide = self._slide()
        self._bar(slide)
        self._title(slide, title)
        self._panel(slide, 0.45, 0.72, 4.35, 4.55)
        self._panel(slide, 5.2, 0.72, 4.35, 4.55)
        self._bullets(slide, left, top=1.0, width=3.95)
        box = slide.shapes.add_textbox(Inches(5.45), Inches(1.0), Inches(3.95), Inches(4.0))
        tf = box.text_frame
        tf.word_wrap = True
        for i, text in enumerate(right):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.font.name = FONT
            p.font.size = Pt(15)
            p.font.color.rgb = DARK
            p.space_after = Pt(8)
        self._footer(slide)

    def image(self, title: str, path: Path, caption: str = "", full: bool = False) -> None:
        slide = self._slide()
        self._bar(slide)
        self._title(slide, title, size=22)
        self._panel(slide, 0.45, 0.68, 9.1, 4.62)
        if path.exists():
            w = Inches(9.0) if full else Inches(8.6)
            slide.shapes.add_picture(str(path), Inches(0.5), Inches(0.78), width=w)
        if caption:
            self._caption(slide, caption)
        self._footer(slide)

    def table(self, title: str, headers: list[str], rows: list[list[str]], highlight: str = "XGBoost") -> None:
        slide = self._slide()
        self._bar(slide)
        self._title(slide, title, size=22)
        self._panel(slide, 0.45, 0.72, 9.1, 4.2)
        nrows, ncols = len(rows) + 1, len(headers)
        shape = slide.shapes.add_table(nrows, ncols, Inches(0.55), Inches(1.05), Inches(8.9), Inches(0.44 * nrows))
        tbl = shape.table
        for j, h in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = h
            for p in c.text_frame.paragraphs:
                p.font.name = FONT
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.color.rgb = WHITE
            c.fill.solid()
            c.fill.fore_color.rgb = NAVY
        for i, row in enumerate(rows, 1):
            for j, val in enumerate(row):
                c = tbl.cell(i, j)
                c.text = val
                for p in c.text_frame.paragraphs:
                    p.font.name = FONT
                    p.font.size = Pt(11)
                    p.font.color.rgb = DARK
                if row[0] == highlight:
                    c.fill.solid()
                    c.fill.fore_color.rgb = RGBColor(0xCC, 0xFB, 0xF1)
        self._footer(slide)

    def cards(self, title: str, cards: list[tuple[str, str]]) -> None:
        """Four compact finding cards."""
        slide = self._slide()
        self._bar(slide)
        self._title(slide, title)
        colors = [BLUE, ORANGE, TEAL, RGBColor(0x7C, 0x3A, 0xED)]
        for i, (head, body) in enumerate(cards):
            col, row = i % 2, i // 2
            x = 0.55 + col * 4.55
            y = 1.25 + row * 2.05
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.25), Inches(1.85)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = colors[i]
            card.line.width = Pt(2)
            tf = card.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = head
            p.font.name = FONT
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = colors[i]
            p2 = tf.add_paragraph()
            p2.text = body
            p2.font.name = FONT
            p2.font.size = Pt(12)
            p2.font.color.rgb = DARK
        self._footer(slide)

    def appendix_divider(self) -> None:
        slide = self._slide()
        self._bar(slide, TEAL, 5.625)
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.2), Inches(8.8), Inches(1.2))
        p = tb.text_frame.paragraphs[0]
        p.text = "Appendix"
        p.font.name = FONT
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(3.2), Inches(8.8), Inches(0.5))
        sp = sub.text_frame.paragraphs[0]
        sp.text = "Detailed EDA charts & supporting visuals"
        sp.font.name = FONT
        sp.font.size = Pt(16)
        sp.font.color.rgb = BLUE_LIGHT
        self._footer(slide)

    def appendix_image(self, label: str, path: Path) -> None:
        slide = self._slide()
        self._bar(slide, TEAL, 0.06)
        tag = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(0.18), Inches(1.3), Inches(0.3)
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = TEAL
        tag.line.fill.background()
        tp = tag.text_frame.paragraphs[0]
        tp.text = "APPENDIX"
        tp.font.name = FONT
        tp.font.size = Pt(9)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
        tp.alignment = PP_ALIGN.CENTER

        self._title(slide, label, top=0.22, size=20)
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(0.55), Inches(0.85), width=Inches(8.9))
        self._footer(slide)

    def closing(self) -> None:
        slide = self._slide()
        block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        block.fill.solid()
        block.fill.fore_color.rgb = NAVY
        block.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(2.1), Inches(8.8), Inches(1.0))
        p = tb.text_frame.paragraphs[0]
        p.text = "Thank you"
        p.font.name = FONT
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = WHITE
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(3.1), Inches(8.8), Inches(0.5))
        sp = sub.text_frame.paragraphs[0]
        sp.text = "Questions?  ·  Group 1  ·  Capstone Iteration 1"
        sp.font.name = FONT
        sp.font.size = Pt(16)
        sp.font.color.rgb = BLUE_LIGHT
        sp.alignment = PP_ALIGN.CENTER
        self._footer(slide)

    def save(self) -> None:
        self.prs.save(str(OUT))
        self.prs.save(str(OUT_PROJECT))
        print(f"Saved {len(self.prs.slides)} slides → {OUT}")


def build_main(d: Deck) -> None:
    d.title_slide(
        "Energy Demand Forecasting — NY Region",
        [
            "Group 1  ·  Capstone Iteration 1",
            "Xuan Toan Doan · MD Shahriar Rashid · Sangaranarayanan SV · Mohammad Shakeel · Sai Susanth",
            "Supervisors: Course Instructor · Industry Mentor (Ms. Saba)",
        ],
    )

    d.content(
        "Iteration 1 — What we delivered",
        [
            "EDA on NY electricity demand (2015–2026)",
            "Frozen train/test dataset via DVC + DagsHub",
            "4 forecasting models: XGBoost, LightGBM, Prophet, SARIMAX",
            "Streamlit dashboard for comparison & demo",
        ],
    )

    d.two_col(
        "Data & preprocessing",
        [
            "EIA Grid Monitor — hourly demand (MW)",
            "Open-Meteo — weather features",
            "Holidays, GDP, population",
        ],
        [
            "Merge on timestamp",
            "Engineer: DOW, season, lags",
            "80/20 chronological split",
            "→ features_selected_train/test.csv",
        ],
    )

    d.content(
        "EDA — demand patterns",
        [
            "Summer & winter peaks; spring/fall troughs",
            "Daily: off-peak 3–5 AM · peak 5–7 PM",
            "Weekdays stable · weekends drop (Sun lowest)",
            "Holidays shift consumption measurably",
        ],
    )

    d.content(
        "EDA — weather drivers",
        [
            "V-shape: demand high at cold & hot temps, lowest ~10–18°C",
            "Snowfall days keep demand elevated",
            "Wind speed shows little marginal effect",
            "Time features (season, hour, DOW) dominate correlations",
        ],
    )

    d.cards(
        "Key EDA findings",
        [
            ("Time features matter", "Season, month, DOW & hour drive most variance"),
            ("V-shape temperature", "Non-linear comfort-band effect on demand"),
            ("Snowfall signal", "Higher baseline demand on snow days"),
            ("Holiday effects", "Predictable but smaller deviations"),
        ],
    )

    # ── Pipeline ─────────────────────────────────────────────────────────────
    d.section("Data pipeline & experiments", "Frozen data · tracked experiments")

    d.image("Architecture", ASSETS / "architecture.png")

    d.content(
        "Frozen dataset (DVC + DagsHub)",
        [
            "DVC tracks hashed CSVs → DagsHub S3 remote",
            "Team pulls via pull_data.py + auth token",
            "Same snapshot = fair model comparison",
            "⚠ Mirroring break blocked full dvc push (Iteration 2 fix)",
        ],
    )

    d.content(
        "Experiment setup",
        [
            "3-day ahead forecast · MLflow on DagsHub",
            "Metrics: MAE, RMSE, MAPE on held-out test",
            "Outputs → reports/*_predictions.csv → Streamlit",
        ],
    )

    # ── Models ───────────────────────────────────────────────────────────────
    d.section("Model results", "Identical test set · 3-day horizon")

    d.image(
        "Model comparison",
        ASSETS / "model_comparison.png",
        "XGBoost lowest avg RMSE; tree models beat statistical baselines.",
    )

    d.table(
        "Test-set metrics (3-day average)",
        ["Model", "Avg MAE", "Avg RMSE", "Avg MAPE"],
        [
            ["XGBoost", "21,037 MW", "29,617 MW", "5.93%"],
            ["LightGBM", "20,877 MW", "29,799 MW", "5.79%"],
            ["Prophet", "23,175 MW", "32,048 MW", "6.27%"],
            ["SARIMAX", "83,515 MW", "96,838 MW", "21.79%"],
        ],
        highlight="XGBoost",
    )

    d.content(
        "Best per model",
        [
            "XGBoost — lowest avg RMSE (29,617 MW) · best overall",
            "LightGBM — lowest avg MAPE (5.79%) · close second on RMSE",
            "Prophet — solid seasonal baseline (32,048 MW RMSE)",
            "SARIMAX — poor OOS; in-sample metrics misleading",
            "LSTM — planned for Iteration 2",
        ],
    )

    # ── UI ───────────────────────────────────────────────────────────────────
    d.section("Streamlit demo", "streamlit run app/streamlit_app.py")

    d.two_col(
        "Dashboard",
        [
            "Model picker (sorted by RMSE)",
            "3-day forecast table + metrics",
            "Actual vs predicted chart",
        ],
        [
            "Compare-all panel",
            "Full test-set leaderboard",
            "Overlay all models on one chart",
        ],
    )

    d.image(
        "Streamlit app",
        ASSETS / "streamlit_app.png",
        "Live dashboard — model picker, 3-day forecast, compare-all panel.",
        full=True,
    )
    d.image("UI — model overlay", ASSETS / "ui_all_models.png")

    # ── Feature importance ───────────────────────────────────────────────────
    d.section("EDA ↔ model validation")

    d.image(
        "XGBoost feature importance",
        ASSETS / "xgboost_feature_importance.png",
        "demand_roll_mean_7d, apparent_temperature (#2), day_of_week.",
    )

    d.image(
        "EDA alignment — all models",
        ASSETS / "eda_alignment_matrix.png",
        "How each model's top features map back to EDA findings.",
    )

    d.image(
        "EDA alignment detail",
        ASSETS / "eda_alignment_all_models.png",
    )

    d.content(
        "Challenges & Iteration 2",
        [
            "DagsHub / DVC / Git sync · mirroring break on push",
            "Test |temp − 14°C| & snowfall flag features",
            "Finalize LSTM · tune SARIMAX · feature selection",
        ],
    )

    d.closing()


def build_appendix(d: Deck) -> None:
    d.appendix_divider()
    model_names = {"xgboost": "XGBoost", "lightgbm": "LightGBM", "prophet": "Prophet", "sarimax": "SARIMAX"}
    for key, label in model_names.items():
        path = ASSETS / f"eda_alignment_{key}.png"
        if path.exists():
            d.appendix_image(f"EDA alignment — {label}", path)
    if not DRAFT_IMAGES.exists():
        return
    for slide_num in sorted(APPENDIX_LABELS):
        label = APPENDIX_LABELS[slide_num]
        images = sorted(DRAFT_IMAGES.glob(f"slide{slide_num:02d}_*"))
        if not images:
            continue
        main = max(images, key=lambda p: p.stat().st_size)
        d.appendix_image(label, main)


def main() -> None:
    import extract_draft_images
    import generate_assets

    extract_draft_images.main()
    generate_assets.main()
    try:
        import capture_streamlit

        capture_streamlit.capture()
    except Exception as exc:
        print(f"Warning: Streamlit screenshot skipped ({exc})")

    d = Deck()
    build_main(d)
    build_appendix(d)
    d.save()


if __name__ == "__main__":
    main()
